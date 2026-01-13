"""
Модуль для извлечения текста из различных типов документов
"""
import io
import logging
from typing import Optional

logger = logging.getLogger("text_extractors")


def extract_text_from_pdf(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из PDF файла
    
    Args:
        file_content: Байты PDF файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        import PyPDF2
        
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text_parts = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        
        return '\n\n'.join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из PDF: %s", e)
        return None


def extract_text_from_docx(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из Word документа (.docx)
    
    Args:
        file_content: Байты DOCX файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        from docx import Document
        
        doc_file = io.BytesIO(file_content)
        doc = Document(doc_file)
        
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        return '\n'.join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из DOCX: %s", e)
        return None


def extract_text_from_doc(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из Word документа (.doc) - старый формат
    
    Args:
        file_content: Байты DOC файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        import textract
        
        # Сохраняем во временный файл, так как textract работает с файлами
        import tempfile
        import os
        
        with tempfile.NamedTemporaryFile(delete=False, suffix='.doc') as tmp_file:
            tmp_file.write(file_content)
            tmp_file_path = tmp_file.name
        
        try:
            # Извлекаем текст с помощью textract
            text = textract.process(tmp_file_path).decode('utf-8')
            return text if text.strip() else None
        finally:
            # Удаляем временный файл
            if os.path.exists(tmp_file_path):
                os.unlink(tmp_file_path)
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из DOC: %s", e)
        return None


def extract_text_from_pptx(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из PowerPoint презентации (.pptx)
    
    Args:
        file_content: Байты PPTX файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        from pptx import Presentation
        
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            if slide_texts:
                text_parts.append(f"Слайд {slide_num}:\n" + "\n".join(slide_texts))
        
        return '\n\n'.join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из PPTX: %s", e)
        return None


def extract_text_from_txt(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из TXT файла
    
    Args:
        file_content: Байты TXT файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        # Пробуем разные кодировки
        for encoding in ['utf-8', 'cp1251', 'latin-1']:
            try:
                return file_content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из TXT: %s", e)
        return None


def extract_text_from_code(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из исходных кодов (py/js/ts/c/cpp/etc.)
    Использует ту же логику, что и для TXT.
    """
    return extract_text_from_txt(file_content)


def extract_text_from_rtf(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из RTF файла
    
    Args:
        file_content: Байты RTF файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(file_content.decode('utf-8', errors='ignore'))
        return text if text.strip() else None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из RTF: %s", e)
        return None


def extract_text_from_xlsx(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из Excel файла (.xlsx)
    
    Args:
        file_content: Байты XLSX файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        from openpyxl import load_workbook
        xlsx_file = io.BytesIO(file_content)
        wb = load_workbook(xlsx_file, data_only=True)
        
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts = []
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    sheet_texts.append(row_text)
            if sheet_texts:
                text_parts.append(f"Лист '{sheet_name}':\n" + "\n".join(sheet_texts))
        
        return '\n\n'.join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из XLSX: %s", e)
        return None


def extract_text_from_xls(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из Excel файла (.xls)
    
    Args:
        file_content: Байты XLS файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        import xlrd
        xls_file = io.BytesIO(file_content)
        wb = xlrd.open_workbook(file_contents=xls_file.read())
        
        text_parts = []
        for sheet in wb.sheets():
            sheet_texts = []
            for row_idx in range(sheet.nrows):
                row = sheet.row_values(row_idx)
                row_text = ' | '.join([str(cell) if cell else '' for cell in row])
                if row_text.strip():
                    sheet_texts.append(row_text)
            if sheet_texts:
                text_parts.append(f"Лист '{sheet.name}':\n" + "\n".join(sheet_texts))
        
        return '\n\n'.join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из XLS: %s", e)
        return None


def extract_text_from_csv(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из CSV файла
    
    Args:
        file_content: Байты CSV файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        import csv
        csv_file = io.BytesIO(file_content)
        
        # Пробуем разные кодировки
        for encoding in ['utf-8', 'cp1251', 'latin-1']:
            try:
                csv_file.seek(0)
                text = csv_file.read().decode(encoding)
                csv_file.seek(0)
                reader = csv.reader(io.StringIO(text))
                rows = [' | '.join(row) for row in reader if any(cell.strip() for cell in row)]
                return '\n'.join(rows) if rows else None
            except (UnicodeDecodeError, Exception):
                continue
        return None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из CSV: %s", e)
        return None


def extract_text_from_md(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из Markdown файла
    
    Args:
        file_content: Байты MD файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        # Пробуем разные кодировки
        for encoding in ['utf-8', 'cp1251', 'latin-1']:
            try:
                return file_content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из MD: %s", e)
        return None


def extract_text_from_html(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из HTML файла
    
    Args:
        file_content: Байты HTML файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        from bs4 import BeautifulSoup
        
        # Пробуем разные кодировки
        for encoding in ['utf-8', 'cp1251', 'latin-1']:
            try:
                html = file_content.decode(encoding)
                # Используем lxml если доступен, иначе html.parser
                soup = BeautifulSoup(html, 'lxml')
                # Удаляем скрипты и стили
                for script in soup(["script", "style"]):
                    script.decompose()
                text = soup.get_text(separator='\n', strip=True)
                if text and text.strip():
                    return text
                # fallback: вернуть сырой декодированный html, если парсинг дал пусто
                if html.strip():
                    return html
            except (UnicodeDecodeError, Exception) as e:
                logger.warning("Ошибка при извлечении текста из HTML: %s", e)
                continue
        return None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из HTML: %s", e)
        return None


def extract_text_from_odt(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из OpenDocument Text (.odt)
    
    Args:
        file_content: Байты ODT файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        from odf import text, teletype
        from odf.opendocument import load
        
        odt_file = io.BytesIO(file_content)
        doc = load(odt_file)
        
        text_parts = []
        for paragraph in doc.getElementsByType(text.P):
            para_text = teletype.extractText(paragraph)
            if para_text.strip():
                text_parts.append(para_text)
        
        return '\n'.join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из ODT: %s", e)
        return None


def extract_text_from_odp(file_content: bytes) -> Optional[str]:
    """
    Извлекает текст из OpenDocument Presentation (.odp)
    
    Args:
        file_content: Байты ODP файла
    
    Returns:
        Извлеченный текст или None в случае ошибки
    """
    try:
        from odf import text, teletype, draw
        from odf.opendocument import load
        
        odp_file = io.BytesIO(file_content)
        doc = load(odp_file)
        
        text_parts = []
        for slide_num, page in enumerate(doc.getElementsByType(draw.Page), 1):
            slide_texts = []
            for paragraph in page.getElementsByType(text.P):
                para_text = teletype.extractText(paragraph)
                if para_text.strip():
                    slide_texts.append(para_text)
            if slide_texts:
                text_parts.append(f"Слайд {slide_num}:\n" + "\n".join(slide_texts))
        
        return '\n\n'.join(text_parts) if text_parts else None
    except Exception as e:
        logger.warning("Ошибка при извлечении текста из ODP: %s", e)
        return None


def extract_text_from_file(file_content: bytes, file_extension: str) -> Optional[str]:
    """
    Универсальная функция для извлечения текста из файла по его расширению
    
    Args:
        file_content: Байты файла
        file_extension: Расширение файла (например, 'pdf', 'docx', 'pptx')
    
    Returns:
        Извлеченный текст или None в случае ошибки или неподдерживаемого формата
    """
    extension = file_extension.lower().lstrip('.')
    
    extractors = {
        'pdf': extract_text_from_pdf,
        'docx': extract_text_from_docx,
        'doc': extract_text_from_doc,
        'pptx': extract_text_from_pptx,
        'txt': extract_text_from_txt,
        'py': extract_text_from_code,
        'js': extract_text_from_code,
        'ts': extract_text_from_code,
        'tsx': extract_text_from_code,
        'c': extract_text_from_code,
        'cpp': extract_text_from_code,
        'h': extract_text_from_code,
        'hpp': extract_text_from_code,
        'java': extract_text_from_code,
        'go': extract_text_from_code,
        'rtf': extract_text_from_rtf,
        'xlsx': extract_text_from_xlsx,
        'xls': extract_text_from_xls,
        'csv': extract_text_from_csv,
        'md': extract_text_from_md,
        'html': extract_text_from_html,
        'htm': extract_text_from_html,
        'odt': extract_text_from_odt,
        'odp': extract_text_from_odp,
    }
    
    extractor = extractors.get(extension)
    if extractor:
        return extractor(file_content)
    else:
        logger.info("Неподдерживаемый формат файла: %s", extension)
        return None


async def extract_text_from_url(url: str, max_bytes: int = 2 * 1024 * 1024) -> Optional[str]:
    """
    Загружает страницу по URL и пытается извлечь из нее текст.
    Ограничивает скачиваемый размер max_bytes, чтобы избежать перегрузки.
    """
    try:
        import aiohttp
    except Exception as e:
        logger.warning("aiohttp недоступен для запроса %s: %s", url, e)
        return None

    try:
        timeout = aiohttp.ClientTimeout(total=15)
        async with aiohttp.ClientSession(timeout=timeout) as session:
            async with session.get(url, allow_redirects=True) as response:
                if response.status != 200:
                    logger.warning("Неуспешный HTTP статус %s для %s", response.status, url)
                    return None

                content_type = (response.headers.get("Content-Type") or "").lower()
                raw = await response.content.read(max_bytes)

                # HTML или любые text/* форматы обрабатываем через HTML-парсер + fallback в текст
                if (
                    "text/html" in content_type
                    or "application/xhtml" in content_type
                    or content_type.startswith("text/")
                    or not content_type
                ):
                    text = extract_text_from_html(raw)
                    if text:
                        return text
                    # Fallback: пробуем напрямую декодировать как текст
                    for encoding in ["utf-8", "cp1251", "latin-1"]:
                        try:
                            return raw.decode(encoding)
                        except UnicodeDecodeError:
                            continue
                # Последний шанс — попытаться вытащить текст как HTML
                return extract_text_from_html(raw)
    except Exception as e:
        logger.warning("Ошибка при извлечении текста по URL %s: %s", url, e)
        return None